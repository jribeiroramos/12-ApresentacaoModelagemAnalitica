from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    #################################################
    # SLIDE 1 – Título & Apresentação
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Análise Multidimensional: Fundamentos e Aplicações"
    subtitle.text = (
        "Como organizar e explorar dados para melhores insights\n"
        "Nome do palestrante/time - Data"
    )

    #################################################
    # SLIDE 2 – Agenda
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Agenda"
    content.text = (
        "1. O que é Análise Multidimensional?\n"
        "2. Conceitos Fundamentais: Fatos & Dimensões\n"
        "3. Hierarquias e Cubos OLAP\n"
        "4. Modelos de Dados (Star e Snowflake)\n"
        "5. ETL e Data Warehouse\n"
        "6. KPIs e Métricas\n"
        "7. Benefícios e Aplicações\n"
        "8. Diferentes Níveis de Usuários\n"
        "9. Roteiro Geral de Entrega de Análise\n"
        "10. Exemplo Prático (Aumento de Cobertura de Municípios)\n"
        "11. Modelagem em Estrela (Detalhes)\n"
        "12. Conclusão e Perguntas\n"
        "13. Descrição da Base Fictícia\n"
        "14. Perguntas (Operacional)\n"
        "15. Perguntas (Gerencial)\n"
        "16. Perguntas (Estratégico)"
    )

    #################################################
    # SLIDE 3 – O que é Análise Multidimensional?
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "O que é Análise Multidimensional?"
    slide.placeholders[1].text = (
        "Definição e Objetivo:\n"
        "- Organizar dados sob diferentes perspectivas.\n"
        "- Essencial para BI (Business Intelligence) e tomada de decisão.\n"
        "- Responde a 'Quando?', 'Onde?', 'Por quê?', 'Quem?', 'Como?'."
    )

    #################################################
    # SLIDE 4 – Conceitos Fundamentais: Fatos
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conceitos Fundamentais: Fatos"
    slide.placeholders[1].text = (
        "O que são Fatos:\n"
        "- Medidas numéricas ou métricas analisadas.\n"
        "- Exemplos: Valor total das vendas, quantidade de itens,\n"
        "  descontos, custos, lucros.\n"
        "- Ficam no centro da análise (Tabela Fato)."
    )

    #################################################
    # SLIDE 5 – Conceitos Fundamentais: Dimensões
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conceitos Fundamentais: Dimensões"
    slide.placeholders[1].text = (
        "O que são Dimensões:\n"
        "- Fornecem contexto aos Fatos.\n"
        "- Respondem: 'Quem?', 'Onde?', 'Quando?', 'Qual produto?'.\n"
        "- Exemplos: Tempo, Produto, Região, Cliente."
    )

    #################################################
    # SLIDE 6 – Hierarquias e Cubos OLAP
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Hierarquias e Cubos OLAP"
    slide.placeholders[1].text = (
        "Hierarquias:\n"
        "- Ex. (Dimensão Tempo): Dia → Mês → Trimestre → Ano.\n"
        "- Ex. (Dimensão Geografia): Cidade → Estado → País.\n\n"
        "Cubos OLAP:\n"
        "- Estruturas multidimensionais para análise rápida.\n"
        "- Permitem drill-down (detalhar) e roll-up (agregar)."
    )

    #################################################
    # SLIDE 7 – Modelos de Dados: Star Schema vs Snowflake
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Modelos de Dados: Star vs Snowflake"
    slide.placeholders[1].text = (
        "Star Schema:\n"
        "- Tabela de Fatos no centro + Dimensões ligadas diretamente.\n"
        "- Mais simples e intuitivo.\n\n"
        "Snowflake Schema:\n"
        "- Dimensões podem ser normalizadas em tabelas adicionais.\n"
        "- Mais complexo, mas pode ajudar no desempenho."
    )

    #################################################
    # SLIDE 8 – ETL e Data Warehouse
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "ETL e Data Warehouse"
    slide.placeholders[1].text = (
        "ETL (Extract, Transform, Load):\n"
        "1. Extrair dados (ERP, CRM, planilhas, etc.).\n"
        "2. Transformar (limpeza, unificação, padronização).\n"
        "3. Carregar em Data Warehouse ou Data Lakehouse.\n\n"
        "Data Warehouse:\n"
        "- Repositório central, confiável para análises."
    )

    #################################################
    # SLIDE 9 – KPIs, Benefícios e Aplicações
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "KPIs, Benefícios e Aplicações"
    slide.placeholders[1].text = (
        "KPIs (Indicadores-Chave):\n"
        "- Ex.: Receita, Margem de Lucro, Ticket Médio, NPS.\n\n"
        "Benefícios:\n"
        "- Visualização em diferentes ângulos.\n"
        "- Insights rápidos (drill-down e roll-up).\n"
        "- Consistência de informação (única 'versão da verdade').\n"
        "- Suporte à tomada de decisão baseada em dados."
    )

    #################################################
    # SLIDE 10 – Benefícios Complementares
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Benefícios Complementares"
    slide.placeholders[1].text = (
        "- Aumenta a transparência e confiança nos dados.\n"
        "- Facilita integração entre áreas (Finanças, Vendas, etc.).\n"
        "- Permite monitoramento de metas e resultados contínuos."
    )

    #################################################
    # SLIDE 11 – Diferentes Níveis de Usuários
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Visão para Diferentes Níveis de Usuários"
    slide.placeholders[1].text = (
        "• **Operacional**: foco no presente e operação diária.\n"
        "• **Gerencial**: análise do presente vs. metas e passado recente.\n"
        "• **Estratégico**: avaliação de tendências passadas e projeções futuras."
    )

    #################################################
    # SLIDE 12 – Tabela Comparativa (Operacional, Gerencial, Estratégico)
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "Comparativo: Operacional, Gerencial e Estratégico"

    from pptx.util import Inches, Pt
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(3.0)

    rows = 4
    cols = 5
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Cabeçalho
    table.cell(0, 0).text = "Nível"
    table.cell(0, 1).text = "Perspectiva\nPrincipal"
    table.cell(0, 2).text = "Periodicidade"
    table.cell(0, 3).text = "Ferramentas"
    table.cell(0, 4).text = "Exemplos de KPIs"

    # Operacional
    table.cell(1, 0).text = "Operacional"
    table.cell(1, 1).text = "Presente (execução diária),\nPassado imediato"
    table.cell(1, 2).text = "Minutos,\nHoras,\nDiário"
    table.cell(1, 3).text = "Dashboards\nem tempo real\n(Grafana,\nKibana,\nPower BI)"
    table.cell(1, 4).text = "Chamados/dia,\nProdução/hora,\nVendas diárias"

    # Gerencial
    table.cell(2, 0).text = "Gerencial"
    table.cell(2, 1).text = "Presente (vs. metas),\nPassado recente"
    table.cell(2, 2).text = "Diário,\nSemanal,\nMensal"
    table.cell(2, 3).text = "BI tradicional\n(Power BI,\nTableau,\nQlik)"
    table.cell(2, 4).text = "Vendas x Meta,\nLucro,\nCobertura"

    # Estratégico
    table.cell(3, 0).text = "Estratégico"
    table.cell(3, 1).text = "Passado (tendências),\nFuturo (projeções)"
    table.cell(3, 2).text = "Mensal,\nTrimestral,\nAnual"
    table.cell(3, 3).text = "BI avançado +\nPlanejamento\n(SAP, Anaplan,\nData Science)"
    table.cell(3, 4).text = "Cenários what-if,\nROI,\nProjeções"

    # Ajuste de fonte
    for r in range(rows):
        for c in range(cols):
            for paragraph in table.cell(r, c).text_frame.paragraphs:
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(11)

    #################################################
    # SLIDE 13 – Roteiro Geral de Entrega de Análise
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Roteiro Geral de Entrega de Análise"
    slide.placeholders[1].text = (
        "Quando recebemos uma 'encomenda' de análise:\n\n"
        "1. Definir o público e objetivo.\n"
        "2. Mapear perspectiva (passado, presente, futuro).\n"
        "3. Selecionar KPIs e dimensões.\n"
        "4. Verificar fontes e qualidade dos dados.\n"
        "5. Modelar e preparar ambiente de análise.\n"
        "6. Escolher a ferramenta de visualização.\n"
        "7. Validar e ajustar.\n"
        "8. Implantar e manter."
    )

    #################################################
    # SLIDE 14 – Passos 1 e 2 (genérico)
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Passo 1 e Passo 2"
    slide.placeholders[1].text = (
        "**1. Definição da Demanda e Público-Alvo**\n"
        "- Entender por que e para quem a análise é solicitada.\n"
        "- Classificar se é nível operacional, gerencial ou estratégico.\n\n"
        "**2. Mapeamento da Perspectiva**\n"
        "- Passado: diagnosticar o que ocorreu.\n"
        "- Presente: monitorar desempenho vs. metas.\n"
        "- Futuro: projeções e cenários what-if."
    )

    #################################################
    # SLIDE 15 – Passos 3 e 4 (genérico)
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Passo 3 e Passo 4"
    slide.placeholders[1].text = (
        "**3. Identificação de KPIs, Métricas e Dimensões**\n"
        "- Quais KPIs respondem às perguntas do negócio?\n"
        "- Quais dimensões (Tempo, Região, Produto etc.) serão usadas?\n\n"
        "**4. Coleta e Qualidade dos Dados**\n"
        "- Fontes (ERP, CRM, planilhas).\n"
        "- Limpeza de dados e padronização.\n"
        "- Definir periodicidade (diário, semanal etc.)."
    )

    #################################################
    # SLIDE 16 – Passos 5 e 6 (genérico)
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Passo 5 e Passo 6"
    slide.placeholders[1].text = (
        "**5. Modelagem e Preparação Analítica**\n"
        "- Estruturar fatos/dimensões (Star, Snowflake).\n"
        "- Configurar processos de ETL.\n"
        "- Definir medidas e cálculos.\n\n"
        "**6. Escolha da Ferramenta**\n"
        "- BI (Power BI, Tableau, Qlik) ou relatórios.\n"
        "- Frequência de atualização.\n"
        "- Nível de interatividade."
    )

    #################################################
    # SLIDE 17 – Passos 7 e 8 (genérico)
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Passo 7 e Passo 8"
    slide.placeholders[1].text = (
        "**7. Validação e Iterações**\n"
        "- Apresentar protótipo aos stakeholders.\n"
        "- Ajustar métricas, layout.\n\n"
        "**8. Implantação e Manutenção**\n"
        "- Publicar relatório/dashboard.\n"
        "- Treinamento e documentação.\n"
        "- Evoluir conforme surgem novas necessidades."
    )

    #################################################
    # SLIDE 18 – Exemplo Prático: Introdução
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Exemplo Prático"
    slide.placeholders[1].text = (
        "Aumento de Cobertura de Municípios (de X para Y)\n\n"
        "• Objetivo: Atingir a meta de cobertura em 1 ano.\n"
        "• Público-Alvo: Gerencial e/ou Estratégico.\n"
        "• Perguntas principais:\n"
        "  - Quantos municípios cobertos no passado?\n"
        "  - Quantos estão cobertos agora?\n"
        "  - Quando atingiremos Y se mantido o ritmo atual?\n"
        "  - Precisamos de estratégias extras?"
    )

    #################################################
    # SLIDE 19 – Passos 1, 2 e 3 no Exemplo de Cobertura
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Exemplo: Passos 1, 2 e 3"
    slide.placeholders[1].text = (
        "**1. Definição do Objetivo e Público-Alvo**\n"
        "- Aumentar de X para Y municípios em 1 ano.\n"
        "- Usuários Gerenciais (acompanhamento) e Estratégicos (decisão).\n\n"
        "**2. Mapeamento (Passado, Presente, Futuro)**\n"
        "- Passado: ritmo histórico de expansão.\n"
        "- Presente: cobertos hoje vs. meta.\n"
        "- Futuro: projeções 'what if'.\n\n"
        "**3. KPIs, Métricas e Dimensões**\n"
        "- KPI Principal: Contagem de municípios cobertos.\n"
        "- Dimensões: Tempo (mês, ano), Geografia (UF, Região).\n"
        "- Ex.: Adições mensais, % coberto vs. total."
    )

    #################################################
    # SLIDE 20 – Passos 4 e 5 no Exemplo de Cobertura
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Exemplo: Passos 4 e 5"
    slide.placeholders[1].text = (
        "**4. Coleta e Qualidade de Dados**\n"
        "- CRM / ERP que registra onde há cobertura.\n"
        "- Base oficial (IBGE) para nomes/códigos de municípios.\n"
        "- Validar e padronizar nomes duplicados, datas faltantes.\n\n"
        "**5. Modelagem e Preparação Analítica**\n"
        "- Fato_Cobertura (1/0 se coberto no período).\n"
        "- Dim_Tempo (dia, mês, ano), Dim_Municipio (nome, estado).\n"
        "- ETL: unificar dados, gerar tabela de fato e dimensões."
    )

    #################################################
    # SLIDE 21 – Exemplo: Passos 6 e 7
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Exemplo: Passos 6 e 7"
    slide.placeholders[1].text = (
        "**6. Escolha da Ferramenta de Visualização**\n"
        "- Dashboard (Power BI, Tableau, Qlik) + mapa geográfico.\n"
        "- Atualização mensal ou semanal.\n\n"
        "**7. Validação e Iterações**\n"
        "- Criar protótipo inicial.\n"
        "- Validar se os números batem com expectativas.\n"
        "- Ajustar rótulos, filtros, layout."
    )

    #################################################
    # SLIDE 22 – Exemplo: Passo 8 e Produto Final
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Exemplo: Passo 8 e Produto Final"
    slide.placeholders[1].text = (
        "**8. Implantação e Manutenção**\n"
        "- Publicar dashboard no portal de BI.\n"
        "- Definir periodicidade (mensal/semanal).\n"
        "- Treinar usuários para interpretar resultados.\n\n"
        "**Produto Final**:\n"
        "- Dashboard com:\n"
        "  • Evolução histórica de municípios.\n"
        "  • Situação atual vs. meta.\n"
        "  • Projeções para atingir Y (cenários)."
    )

    #################################################
    # SLIDE 23 – Introdução à Modelagem em Estrela
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Modelagem em Estrela (Star Schema)"
    slide.placeholders[1].text = (
        "Detalhes do cenário de Cobertura de Municípios:\n\n"
        "• Tabela de Fato: Fato_Cobertura\n"
        "• Dimensões: Dim_Tempo, Dim_Municipio, Dim_Estrategia\n"
        "• Permite analisar onde (município), quando (tempo)\n"
        "  e como (estratégia) a cobertura foi estabelecida."
    )

    #################################################
    # SLIDE 24 – Visão Geral do Star Schema
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Visão Geral do Star Schema"
    slide.placeholders[1].text = (
        "Fato_Cobertura\n"
        "  - FK para Dim_Tempo (ID_Tempo)\n"
        "  - FK para Dim_Municipio (ID_Munic)\n"
        "  - FK para Dim_Estrategia (ID_Estrategia)\n"
        "  - Atributos: Coberto (1/0), Datas de início/fim\n\n"
        "Dim_Tempo\n"
        "  - Dia, Mês, Ano, etc.\n\n"
        "Dim_Municipio\n"
        "  - Nome, Estado, Região, Código IBGE etc.\n\n"
        "Dim_Estrategia\n"
        "  - Nome, Tipo, Equipe Responsável etc."
    )

    #################################################
    # SLIDE 25 – Tabela de Fato_Cobertura
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Fato_Cobertura: Estrutura e Campos"
    slide.placeholders[1].text = (
        "**Nome:** Fato_Cobertura\n"
        "**Atributos Principais:**\n"
        "• ID_Fato (opcional, PK) ou chave composta\n"
        "• ID_Tempo (FK), ID_Munic (FK), ID_Estrategia (FK)\n"
        "• Coberto (1/0)\n"
        "• Data_Inicio_Cobertura, Data_Fim_Cobertura\n\n"
        "Pode haver uma linha por dia/mês ou apenas no momento\n"
        "de adesão, dependendo da granularidade desejada."
    )

    #################################################
    # SLIDE 26 – Dimensões: Tempo, Município e Estratégia
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dimensões: Tempo, Município e Estratégia"
    slide.placeholders[1].text = (
        "**Dim_Tempo**\n"
        "- ID_Tempo (PK)\n"
        "- Data, Dia, Mês, Ano, etc.\n\n"
        "**Dim_Municipio**\n"
        "- ID_Munic (PK)\n"
        "- Nome_Municipio, Estado, Região, Cod_IBGE...\n\n"
        "**Dim_Estrategia**\n"
        "- ID_Estrategia (PK)\n"
        "- Nome_Estrategia, Tipo_Estrategia, Equipe_Responsavel..."
    )

    #################################################
    # SLIDE 27 – Exemplos de Consultas
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Exemplos de Consultas"
    slide.placeholders[1].text = (
        "1. Quantos municípios cobertos por mês no último ano?\n"
        "   SELECT t.Ano, t.Mes,\n"
        "          COUNT(DISTINCT fc.ID_Munic) AS municipios_cobertos\n"
        "   FROM Fato_Cobertura fc\n"
        "   JOIN Dim_Tempo t ON fc.ID_Tempo = t.ID_Tempo\n"
        "   WHERE t.Ano = 2023 AND fc.Coberto = 1\n"
        "   GROUP BY t.Ano, t.Mes;\n\n"
        "2. Cobertura por Estado e Estratégia:\n"
        "   SELECT m.Estado, e.Nome_Estrategia,\n"
        "          COUNT(DISTINCT fc.ID_Munic) AS total_cobertos\n"
        "   FROM Fato_Cobertura fc\n"
        "   JOIN Dim_Municipio m ON fc.ID_Munic = m.ID_Munic\n"
        "   JOIN Dim_Estrategia e ON fc.ID_Estrategia = e.ID_Estrategia\n"
        "   WHERE fc.Coberto = 1\n"
        "   GROUP BY m.Estado, e.Nome_Estrategia;"
    )

    #################################################
    # SLIDE 28 – Conclusão do Exemplo de Cobertura
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusão do Exemplo de Cobertura"
    slide.placeholders[1].text = (
        "• Modelo em Estrela facilita consultas e análises.\n"
        "• Fato_Cobertura concentra o 'quando' e 'onde'.\n"
        "• Dimensões enriquecem os dados (Tempo, Município, Estratégia).\n"
        "• Permite ver evolução histórica, situação atual e projeções.\n"
        "• Traz clareza para decidir recursos e investimentos na expansão."
    )

    #################################################
    # SLIDE 29 – Perguntas e Discussão
    #################################################
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Perguntas e Discussão"
    slide.placeholders[1].text = (
        "• Dúvidas sobre os conceitos ou o modelo?\n"
        "• Próximos passos no projeto real.\n"
        "• Obrigado pela atenção!"
    )

    #################################################
    # NOVOS SLIDES 30+ – Descrevendo a Base Fictícia e Perguntas SQL
    #################################################

    # SLIDE 30 – Descrição da Base Fictícia
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Descrição da Base Fictícia"
    slide.placeholders[1].text = (
        "• Gera dados de 20 anos (2005 a 2024), com granularidade diária.\n"
        "• 100 municípios fictícios, 5 estratégias diferentes.\n"
        "• Cada município tem uma data de início (e possível data de fim)\n"
        "  de cobertura.\n"
        "• Resulta em ~730k linhas na Fato_Cobertura (diária),\n"
        "  armazenando se está coberto (1) ou não (0) em cada dia."
    )

    # SLIDE 31 – Perguntas (Operacional)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Perguntas (Nível Operacional)"
    slide.placeholders[1].text = (
        "Exemplos de perguntas e SQL:\n\n"
        "1) **Quantos municípios estão cobertos HOJE?**\n"
        "   SELECT COUNT(DISTINCT ID_Munic)\n"
        "   FROM Fato_Cobertura fc\n"
        "   JOIN Dim_Tempo t ON fc.ID_Tempo = t.ID_Tempo\n"
        "   WHERE t.DataCompleta = '2024-12-31' AND fc.Coberto = 1;\n\n"
        "2) **Existe algum município que perdeu cobertura hoje?**\n"
        "   (varia conforme data_fim)."
    )

    # SLIDE 32 – Perguntas (Nível Gerencial)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Perguntas (Nível Gerencial)"
    slide.placeholders[1].text = (
        "Exemplos de perguntas e SQL:\n\n"
        "1) **Cobertura atual vs. meta**:\n"
        "   SELECT (COUNT(DISTINCT fc.ID_Munic)) AS CoberturaAtual,\n"
        "          600 AS Meta -- Exemplo\n"
        "   FROM Fato_Cobertura fc\n"
        "   JOIN Dim_Tempo t ON fc.ID_Tempo = t.ID_Tempo\n"
        "   WHERE t.DataCompleta = '2024-12-31' AND fc.Coberto = 1;\n\n"
        "2) **Cobertura por Região**:\n"
        "   SELECT m.Regiao, COUNT(DISTINCT fc.ID_Munic) AS qtde\n"
        "   FROM Fato_Cobertura fc\n"
        "   JOIN Dim_Tempo t ON fc.ID_Tempo = t.ID_Tempo\n"
        "   JOIN Dim_Municipio m ON fc.ID_Munic = m.ID_Munic\n"
        "   WHERE t.Ano = 2024 AND fc.Coberto = 1\n"
        "   GROUP BY m.Regiao;"
    )

    # SLIDE 33 – Perguntas (Nível Estratégico)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Perguntas (Nível Estratégico)"
    slide.placeholders[1].text = (
        "Exemplos de perguntas e SQL:\n\n"
        "1) **Tendência de Cobertura nos últimos 5 anos**:\n"
        "   SELECT t.Ano, COUNT(DISTINCT fc.ID_Munic) AS cobertos\n"
        "   FROM Fato_Cobertura fc\n"
        "   JOIN Dim_Tempo t ON fc.ID_Tempo = t.ID_Tempo\n"
        "   WHERE t.Ano BETWEEN 2020 AND 2024\n"
        "         AND fc.Coberto = 1\n"
        "   GROUP BY t.Ano;\n\n"
        "2) **Projeções Futuras**:\n"
        "   Aqui, o SQL pode ser combinado com\n"
        "   modelos de previsão (Machine Learning)\n"
        "   para estimar quando atingiremos certa meta.\n"
    )

    # FIM
    prs.save("apresentacao_analise_multidimensional.pptx")
    print("Apresentação gerada com sucesso: apresentacao_analise_multidimensional.pptx")


if __name__ == "__main__":
    create_presentation()

